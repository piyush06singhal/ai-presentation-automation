from pptx.util import Inches, Pt
from typing import List, Tuple, Optional
from app.schemas.models import SlidePlan
from app.services.text_renderer import TextRenderer
from app.services.theme_manager import CorporateTheme

class ContentRenderer:
    """Combines text styles, layouts, and speaker notes to populate structured slide content zones."""

    @staticmethod
    def render_header(slide, slide_plan: SlidePlan):
        """Assembles and writes slide title and sub-objective texts."""
        TextRenderer.add_slide_header(
            slide=slide,
            title_text=slide_plan.title,
            subtitle_text=slide_plan.objective
        )

    @staticmethod
    def render_insights(slide, insights: List[str], bounding_box: Tuple[Inches, Inches, Inches, Inches]):
        """Renders slide insight lists using correct typography constraints."""
        if insights:
            TextRenderer.add_bullet_list(
                slide=slide,
                items=insights,
                bounding_box=bounding_box,
                list_title="Key Analytical Insights"
            )

    @staticmethod
    def render_recommendations(slide, recommendations: List[str], bounding_box: Tuple[Inches, Inches, Inches, Inches]):
        """Renders recommendations in a dedicated box zone."""
        if recommendations:
            TextRenderer.add_bullet_list(
                slide=slide,
                items=recommendations,
                bounding_box=bounding_box,
                list_title="Operational Recommendations"
            )

    @staticmethod
    def apply_speaker_notes(slide, notes_text: str):
        """Writes presenter guide notes into PowerPoint's native slide notes zone."""
        if not notes_text:
            return
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = str(notes_text).strip()
        except Exception:
            # Fallback if notes frame fails to load
            pass
