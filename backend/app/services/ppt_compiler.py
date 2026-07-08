import io
import pandas as pd
from typing import Dict
from pptx import Presentation
from app.schemas.models import StoryboardRequest, BusinessSummary
from app.services.theme_manager import CorporateTheme
from app.services.slide_factory import SlideFactory
from app.services.text_renderer import TextRenderer

class PPTCompiler:
    """Orchestrates layout compilation, theme coloring, slide formatting, and presentation exports."""

    @staticmethod
    def compile_presentation(
        plan: StoryboardRequest,
        summary: BusinessSummary,
        df_collection: Dict[str, pd.DataFrame]
    ) -> bytes:
        """Assembles slides in sequence, applying theme backgrounds, visual elements, and pagination."""
        prs = Presentation()
        
        # 1. Enforce 16:9 Widescreen standard dimensions
        prs.slide_width = CorporateTheme.SLIDE_WIDTH
        prs.slide_height = CorporateTheme.SLIDE_HEIGHT

        total_slides = len(plan.slides)
        
        # 2. Iterate and draw slides
        for slide_idx, slide_plan in enumerate(plan.slides):
            # Use blank layout (index 6 in python-pptx default template) for precise coordinate layout
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            # 3. Apply Background Shade
            background = slide.background
            fill = background.fill
            fill.solid()
            
            # Dark title slides, clean white for content slides
            if slide_plan.template_id == "Title":
                fill.fore_color.rgb = CorporateTheme.BACKGROUND_DARK
            else:
                fill.fore_color.rgb = CorporateTheme.WHITE

            # 4. Draw Slide Contents
            SlideFactory.build_slide(
                slide=slide,
                plan=slide_plan,
                summary=summary,
                df_collection=df_collection
            )

            # 5. Append Footers & Page Numbers (skip on title slide)
            if slide_plan.template_id != "Title":
                TextRenderer.add_slide_footer(
                    slide=slide,
                    slide_num=slide_idx + 1,
                    total_slides=total_slides
                )

        # 6. Save presentation to memory-only stream
        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        
        return output_stream.getvalue()
