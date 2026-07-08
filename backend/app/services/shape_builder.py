from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Tuple
from app.services.theme_manager import CorporateTheme

class ShapeBuilder:
    """Provides methods to draw KPI cards, callout containers, and borders using native shapes."""

    @staticmethod
    def draw_kpi_card(
        slide,
        left: Inches,
        top: Inches,
        width: Inches,
        height: Inches,
        metric_value: str,
        metric_label: str,
        card_color: RGBColor = CorporateTheme.BACKGROUND_LIGHT
    ):
        """Draws a clean, styled numeric card for displaying core business KPIs."""
        # 1. Add background card shape (Rounded Rectangle)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = card_color
        
        # Border style
        shape.line.color.rgb = CorporateTheme.LIGHT_GRAY
        shape.line.width = Pt(1.5)

        # 2. Add text box inside the card zone for rendering
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.1)

        # Large Metric Value
        p = tf.paragraphs[0]
        p.text = str(metric_value).strip()
        p.font.name = CorporateTheme.FONT_TITLE
        p.font.size = Pt(CorporateTheme.SIZE_KPI_VAL)
        p.font.bold = True
        p.font.color.rgb = CorporateTheme.PRIMARY

        # Metric Label/Title below it
        p2 = tf.add_paragraph()
        p2.text = str(metric_label).strip().upper()
        p2.font.name = CorporateTheme.FONT_BODY
        p2.font.size = Pt(CorporateTheme.SIZE_KPI_LABEL)
        p2.font.bold = True
        p2.font.color.rgb = CorporateTheme.SECONDARY
        p2.space_before = Pt(4)

    @staticmethod
    def draw_callout_box(
        slide,
        bounding_box: Tuple[Inches, Inches, Inches, Inches],
        text: str,
        border_color: RGBColor = CorporateTheme.ACCENT,
        fill_color: RGBColor = CorporateTheme.BACKGROUND_LIGHT
    ):
        """Draws a callout box panel with a thick left border."""
        left, top, width, height = bounding_box
        
        # Draw background panel
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = fill_color
        panel.line.fill.background()  # Transparent/no borders

        # Draw left border line indicator
        line_width = Inches(0.08)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, line_width, height
        )
        border.fill.solid()
        border.fill.fore_color.rgb = border_color
        border.line.fill.background()

        # Add text box inside panel
        tx_box = slide.shapes.add_textbox(left + line_width + Inches(0.1), top, width - line_width - Inches(0.2), height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.1)
        
        p = tf.paragraphs[0]
        p.text = str(text).strip()
        p.font.name = CorporateTheme.FONT_BODY
        p.font.size = Pt(CorporateTheme.SIZE_BODY)
        p.font.color.rgb = CorporateTheme.PRIMARY
