import pandas as pd
import traceback
from typing import Dict, List, Tuple, Optional
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.schemas.models import SlidePlan, BusinessSummary, DataTypeEnum, DetectedKPI
from app.services.theme_manager import CorporateTheme
from app.services.layout_manager import LayoutManager
from app.services.text_renderer import TextRenderer
from app.services.content_renderer import ContentRenderer
from app.services.shape_builder import ShapeBuilder
from app.services.chart_builder import ChartBuilder
from app.services.table_builder import TableBuilder

class SlideFactory:
    """Dispatches layout construction tasks to build formatted presentation slides."""

    @staticmethod
    def build_slide(
        slide,
        plan: SlidePlan,
        summary: BusinessSummary,
        df_collection: Dict[str, pd.DataFrame]
    ):
        """Orchestrates header, visual layouts, insights, and speaker notes mapping."""
        template_id = plan.template_id
        
        # 1. Draw Slide Titles & Headers (unless it is the main Title slide)
        if template_id != "Title":
            ContentRenderer.render_header(slide, plan)
            
        # 2. Compile visual template layouts
        if template_id == "Title":
            SlideFactory._build_title_slide(slide, plan)
        elif template_id == "Agenda":
            SlideFactory._build_agenda_slide(slide, plan, summary)
        elif template_id == "KPI Dashboard":
            SlideFactory._build_kpi_dashboard_slide(slide, plan, summary, df_collection)
        elif template_id == "Trend Analysis" or template_id == "Category Comparison":
            SlideFactory._build_chart_slide(slide, plan, summary, df_collection)
        elif template_id == "Summary Table":
            SlideFactory._build_table_slide(slide, plan, summary, df_collection)
        elif template_id == "Recommendations":
            SlideFactory._build_recommendation_slide(slide, plan)
        elif template_id == "Appendix":
            SlideFactory._build_appendix_slide(slide, plan, summary)
        else:
            # Fallback layout: standard executive summary layout (Text Flow)
            content_box = LayoutManager.get_full_content_box()
            ContentRenderer.render_insights(slide, plan.insights, content_box)

        # 3. Add Speaker notes
        ContentRenderer.apply_speaker_notes(slide, plan.speaker_notes)

    @staticmethod
    def _build_title_slide(slide, plan: SlidePlan):
        """Draws a clean, centered title page layout."""
        # 1. Main Title
        left = Inches(1.0)
        top = Inches(2.2)
        width = CorporateTheme.SLIDE_WIDTH - Inches(2.0)
        height = Inches(1.8)
        
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = plan.title.upper()
        p.font.name = CorporateTheme.FONT_TITLE
        p.font.size = Pt(CorporateTheme.SIZE_TITLE_MAIN)
        p.font.bold = True
        p.font.color.rgb = CorporateTheme.PRIMARY
        
        # 2. Objective Subtitle
        p2 = tf.add_paragraph()
        p2.text = plan.objective
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = CorporateTheme.FONT_BODY
        p2.font.size = Pt(CorporateTheme.SIZE_BODY + 4)
        p2.font.color.rgb = CorporateTheme.ACCENT
        p2.space_before = Pt(12)

        # 3. Date / Presenter metadata block at bottom
        meta_left = Inches(2.0)
        meta_top = Inches(5.0)
        meta_width = CorporateTheme.SLIDE_WIDTH - Inches(4.0)
        meta_height = Inches(0.8)
        
        meta_box = slide.shapes.add_textbox(meta_left, meta_top, meta_width, meta_height)
        mtf = meta_box.text_frame
        mtf.word_wrap = True
        
        p3 = mtf.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        p3.text = "Prepared by DeckMate Presentation Automation Platform"
        p3.font.name = CorporateTheme.FONT_BODY
        p3.font.size = Pt(CorporateTheme.SIZE_CAPTION)
        p3.font.color.rgb = CorporateTheme.SECONDARY

    @staticmethod
    def _build_agenda_slide(slide, plan: SlidePlan, summary: BusinessSummary):
        """Renders the agenda slide listing analyzed worksheets as contents."""
        content_box = LayoutManager.get_full_content_box()
        
        agenda_items = []
        for i, s in enumerate(summary.metadata.sheets):
            agenda_items.append(f"Section {i+1}: {s.name} ({s.row_count} Rows Analyzed)")
            
        # Draw Agenda columns list
        ContentRenderer.render_insights(slide, agenda_items, content_box)

    @staticmethod
    def _resolve_columns(
        df: pd.DataFrame,
        worksheet: str,
        summary: BusinessSummary,
        plan_x: Optional[str] = None,
        plan_y: Optional[List[str]] = None
    ) -> Tuple[str, List[str]]:
        """Resolves valid X and Y columns, falling back to sheet schema if plan columns are missing or invalid."""
        # Find sheet metadata
        sheet_meta = None
        for s in summary.metadata.sheets:
            if s.name == worksheet:
                sheet_meta = s
                break
        if not sheet_meta and summary.metadata.sheets:
            sheet_meta = summary.metadata.sheets[0]

        # 1. Resolve X Column
        x_col = None
        if plan_x and plan_x in df.columns:
            x_col = plan_x
        elif sheet_meta:
            # First categorical column
            cat_cols = [c.name for c in sheet_meta.columns if c.datatype in [DataTypeEnum.CATEGORICAL, DataTypeEnum.IDENTIFIER] and c.name in df.columns]
            if cat_cols:
                x_col = cat_cols[0]
        
        if not x_col:
            x_col = df.columns[0] if len(df.columns) > 0 else ""

        # 2. Resolve Y Columns
        y_cols = []
        if plan_y:
            y_cols = [y for y in plan_y if y in df.columns]
            
        if not y_cols and sheet_meta:
            # All numeric columns
            num_cols = [c.name for c in sheet_meta.columns if c.datatype in [DataTypeEnum.NUMERIC, DataTypeEnum.PERCENTAGE, DataTypeEnum.CURRENCY] and c.name in df.columns]
            y_cols = num_cols[:3] # limit to 3

        if not y_cols:
            y_cols = [col for col in df.columns if col != x_col][:2]

        # Prevent duplicate x_col in y_cols
        y_cols = [y for y in y_cols if y != x_col]
        
        return x_col, y_cols

    @staticmethod
    def _draw_top_kpi_row(slide, summary: BusinessSummary, worksheet_name: str, df: Optional[pd.DataFrame] = None):
        """Draws a standardized row of 3 metric cards at the top of the content canvas filtered by worksheet."""
        kpi_row_box = LayoutManager.get_composite_kpi_row_box()
        
        # Filter KPIs specific to the current worksheet
        display_kpis = [k for k in summary.kpis if k.worksheet == worksheet_name]
        
        # If we have less than 3 sheet-specific KPIs, dynamically compute fallback metrics from this sheet's df
        if len(display_kpis) < 3 and df is not None:
            # Dynamically scan numeric columns in df to aggregate
            num_cols = df.select_dtypes(include=['number']).columns.tolist()
            # Remove index columns like "Unnamed: 0"
            num_cols = [c for c in num_cols if c != "Unnamed: 0"]
            
            for col in num_cols:
                if len(display_kpis) >= 3:
                    break
                col_lower = str(col).lower()
                # Skip columns already used in display_kpis
                if any(k.column == col for k in display_kpis):
                    continue
                    
                # Compute average or sum based on column type/name
                col_data = df[col].dropna()
                if not col_data.empty:
                    if "id" in col_lower or "count" in col_lower or "total" in col_lower or "sum" in col_lower:
                        val = float(col_data.sum())
                        label = f"Total {col}"
                    else:
                        val = float(col_data.mean())
                        label = f"Average {col}"
                        
                    unit = "%" if "%" in col_lower or "rate" in col_lower or "margin" in col_lower or "progress" in col_lower else ""
                    if any(k in col_lower for k in ["price", "cost", "revenue", "sales", "spend", "profit"]):
                        unit = "$"
                        
                    display_kpis.append(DetectedKPI(
                        name=label,
                        value=val,
                        column=col,
                        worksheet=worksheet_name,
                        confidence=1.0,
                        unit=unit,
                        description=f"Computed metric for {col}."
                    ))

        # If still less than 3, fallback to generic workbook stats
        if len(display_kpis) < 3:
            if not any(k.name == "Total Rows" for k in display_kpis) and df is not None:
                display_kpis.append(DetectedKPI(
                    name="Total Rows",
                    value=float(len(df)),
                    column="Row Count",
                    worksheet=worksheet_name,
                    confidence=1.0,
                    description="Total rows in worksheet."
                ))
            if len(display_kpis) < 3 and not any(k.name == "Worksheets" for k in display_kpis):
                display_kpis.append(DetectedKPI(
                    name="Worksheets",
                    value=float(len(summary.metadata.sheets)),
                    column="Sheet Count",
                    worksheet="",
                    confidence=1.0,
                    description="Number of parsed tabs."
                ))
            if len(display_kpis) < 3 and not any(k.name == "Health Rating" for k in display_kpis):
                display_kpis.append(DetectedKPI(
                    name="Health Rating",
                    value=float(summary.health.overall_score),
                    column="Health Score",
                    worksheet="",
                    confidence=1.0,
                    unit="%",
                    description="Workbook quality score."
                ))
                
        display_kpis = display_kpis[:3]
        cells = LayoutManager.get_grid_cells(rows=1, cols=3, parent_box=kpi_row_box)
        for idx, cell in enumerate(cells):
            kpi = display_kpis[idx]
            # Format value
            if kpi.unit == "%":
                # If progress/rate is a ratio like 0.78, multiply by 100 for display
                if kpi.value <= 1.0 and kpi.value > 0:
                    val_str = f"{kpi.value * 100:.1f}%"
                else:
                    val_str = f"{kpi.value:.1f}%"
            elif kpi.unit == "$":
                val_str = f"${kpi.value:,.0f}" if kpi.value.is_integer() else f"${kpi.value:,.2f}"
            else:
                val_str = f"{kpi.value:,.0f}" if kpi.value.is_integer() else f"{kpi.value:,.2f}"
                if kpi.unit:
                    val_str += f" {kpi.unit}"
                    
            ShapeBuilder.draw_kpi_card(
                slide=slide,
                left=cell[0],
                top=cell[1],
                width=cell[2],
                height=cell[3],
                metric_value=val_str,
                metric_label=kpi.name
            )

    @staticmethod
    def _build_kpi_dashboard_slide(
        slide,
        plan: SlidePlan,
        summary: BusinessSummary,
        df_collection: Dict[str, pd.DataFrame]
    ):
        """Renders 3 KPI cards arranged at the top, and a full-width summary table at the bottom."""
        # 1. Draw top KPI row
        df = df_collection.get(plan.worksheet)
        SlideFactory._draw_top_kpi_row(slide, summary, plan.worksheet, df)
        
        # 2. Draw bottom summary table
        bottom_box = LayoutManager.get_composite_bottom_full_box()
        df = df_collection.get(plan.worksheet)
        if df is not None:
            # Safe dynamic column resolution
            x_col, y_cols = SlideFactory._resolve_columns(df, plan.worksheet, summary)
            
            try:
                TableBuilder.draw_table(
                    slide=slide,
                    bounding_box=bottom_box,
                    df=df,
                    x_col=x_col,
                    y_cols=y_cols
                )
            except Exception:
                traceback.print_exc()
                ShapeBuilder.draw_callout_box(
                    slide=slide,
                    bounding_box=bottom_box,
                    text="Factual source data summary table could not be loaded into the visual grid.",
                    border_color=CorporateTheme.WARNING
                )
        else:
            # Fallback to insights list if no df exists
            ContentRenderer.render_insights(slide, plan.insights, bottom_box)

    @staticmethod
    def _build_chart_slide(
        slide,
        plan: SlidePlan,
        summary: BusinessSummary,
        df_collection: Dict[str, pd.DataFrame]
    ):
        """Assembles composite slide layouts: 3 KPI cards at top, split chart (left) and insights (right) at bottom."""
        # 1. Draw top KPI row
        df = df_collection.get(plan.worksheet)
        SlideFactory._draw_top_kpi_row(slide, summary, plan.worksheet, df)
        
        # 2. Get composite bottom split zones
        left_box, right_box = LayoutManager.get_composite_bottom_split_boxes()
        
        df = df_collection.get(plan.worksheet)
        chart_success = False
        
        if df is not None:
            x_col, y_cols = SlideFactory._resolve_columns(df, plan.worksheet, summary, plan.x_axis, plan.y_axis)
            
            # Attempt to render native chart
            if plan.chart_type and plan.chart_type.lower() != "none" and y_cols:
                try:
                    ChartBuilder.draw_chart(
                        slide=slide,
                        bounding_box=left_box,
                        chart_type=plan.chart_type,
                        df=df,
                        x_col=x_col,
                        y_cols=y_cols
                    )
                    chart_success = True
                except Exception:
                    traceback.print_exc()
                    chart_success = False

            # Fallback to Summary Table if chart drawing fails or is set to None
            if not chart_success and y_cols:
                try:
                    TableBuilder.draw_table(
                        slide=slide,
                        bounding_box=left_box,
                        df=df,
                        x_col=x_col,
                        y_cols=y_cols
                    )
                except Exception:
                    traceback.print_exc()
                    ShapeBuilder.draw_callout_box(
                        slide=slide,
                        bounding_box=left_box,
                        text="Factual source data columns could not be loaded into visual shape grids.",
                        border_color=CorporateTheme.WARNING
                    )

        # Draw bullet insights in the right column box
        ContentRenderer.render_insights(slide, plan.insights, right_box)

    @staticmethod
    def _build_table_slide(
        slide,
        plan: SlidePlan,
        summary: BusinessSummary,
        df_collection: Dict[str, pd.DataFrame]
    ):
        """Layout slide displaying a split data grid (left) and bullet insights (right) underneath 3 KPI cards."""
        # 1. Draw top KPI row
        df = df_collection.get(plan.worksheet)
        SlideFactory._draw_top_kpi_row(slide, summary, plan.worksheet, df)
        
        # 2. Get composite split zones
        left_box, right_box = LayoutManager.get_composite_bottom_split_boxes()
        df = df_collection.get(plan.worksheet)
        
        if df is not None:
            x_col, y_cols = SlideFactory._resolve_columns(df, plan.worksheet, summary, plan.x_axis, plan.y_axis)
            try:
                TableBuilder.draw_table(
                    slide=slide,
                    bounding_box=left_box,
                    df=df,
                    x_col=x_col,
                    y_cols=y_cols
                )
            except Exception:
                traceback.print_exc()
                ShapeBuilder.draw_callout_box(
                    slide=slide,
                    bounding_box=left_box,
                    text="Failed to display data summary grid.",
                    border_color=CorporateTheme.ERROR
                )
                
        # Insights list
        ContentRenderer.render_insights(slide, plan.insights, right_box)

    @staticmethod
    def _build_recommendation_slide(slide, plan: SlidePlan):
        """Draws horizontal callout cards for each action item."""
        parent_left, parent_top, parent_width, parent_height = LayoutManager.get_full_content_box()
        
        recs = plan.recommendations if plan.recommendations else []
        if not recs:
            # Fallback to standard summary if recommendations list is empty
            ContentRenderer.render_insights(slide, plan.insights, (parent_left, parent_top, parent_width, parent_height))
            return
            
        rows = len(recs)
        cells = LayoutManager.get_grid_cells(rows=rows, cols=1, parent_box=(parent_left, parent_top, parent_width, parent_height))
        
        for idx, cell in enumerate(cells):
            ShapeBuilder.draw_callout_box(
                slide=slide,
                bounding_box=cell,
                text=recs[idx],
                border_color=CorporateTheme.SUCCESS
            )

    @staticmethod
    def _build_appendix_slide(slide, plan: SlidePlan, summary: BusinessSummary):
        """Displays data audit quality indicators and checks status logs."""
        left_box, right_box = LayoutManager.get_split_layout_boxes()
        
        # 1. Left side: Audit issues summary
        health_issues_str = []
        if not summary.health.issues:
            health_issues_str.append("Workbook checks passed. Health status is clean.")
        else:
            for issue in summary.health.issues[:6]:  # Limit to first 6 warnings to avoid overflow
                health_issues_str.append(f"[{issue.severity.value}] Worksheet '{issue.worksheet}': {issue.warning}")
                
        TextRenderer.add_bullet_list(
            slide=slide,
            items=health_issues_str,
            bounding_box=left_box,
            list_title="Workbook Data Quality Audit Log"
        )
        
        # 2. Right side: standard outline notes
        ContentRenderer.render_insights(slide, plan.insights, right_box)
