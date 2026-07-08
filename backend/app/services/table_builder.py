import pandas as pd
from typing import Tuple, List, Dict
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from app.services.theme_manager import CorporateTheme

class TableBuilder:
    """Draws styled database grids inside slides, formatting headers and cell values dynamically."""

    @staticmethod
    def draw_table(
        slide,
        bounding_box: Tuple[Inches, Inches, Inches, Inches],
        df: pd.DataFrame,
        x_col: str,
        y_cols: List[str]
    ):
        """Assembles a native PowerPoint table shape containing the target dataframe segment."""
        left, top, width, height = bounding_box
        
        # Select target dataset slice (limit rows to 10 to fit cleanly inside height boundary)
        df_slice = df[[x_col] + y_cols].head(10)
        
        num_rows = len(df_slice) + 1  # data + 1 header row
        num_cols = len(df_slice.columns)

        # 1. Add Table Shape
        table_shape = slide.shapes.add_table(
            num_rows, num_cols, left, top, width, height
        )
        table = table_shape.table

        # 2. Render Header Row
        for col_idx, col_name in enumerate(df_slice.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name).strip()
            
            # Format text frame
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
            p.font.name = CorporateTheme.FONT_BODY
            p.font.size = Pt(CorporateTheme.SIZE_BODY)
            p.font.bold = True
            p.font.color.rgb = CorporateTheme.WHITE
            
            # Apply Header styling color
            cell.fill.solid()
            cell.fill.fore_color.rgb = CorporateTheme.PRIMARY

        # 3. Render Data Rows
        for row_idx, (_, row_data) in enumerate(df_slice.iterrows()):
            table_row_idx = row_idx + 1
            
            # Alternating shading color
            bg_color = CorporateTheme.BACKGROUND_LIGHT if table_row_idx % 2 == 1 else CorporateTheme.WHITE
            
            for col_idx, col_val in enumerate(row_data):
                cell = table.cell(table_row_idx, col_idx)
                
                # Dynamic formatting check
                if pd.isna(col_val):
                    formatted_val = "-"
                elif col_idx > 0 and isinstance(col_val, (int, float)):
                    # Format as metric currency or percentage
                    col_name = df_slice.columns[col_idx].lower()
                    if "margin" in col_name or "rate" in col_name or "%" in col_name or "growth" in col_name:
                        formatted_val = f"{col_val:.1%}" if col_val <= 1.0 else f"{col_val:.1f}%"
                    elif any(k in col_name for k in ["price", "cost", "revenue", "sales", "spend", "profit"]):
                        formatted_val = f"${col_val:,.2f}"
                    else:
                        formatted_val = f"{col_val:,.0f}" if col_val.is_integer() else f"{col_val:,.2f}"
                else:
                    formatted_val = str(col_val).strip()

                cell.text = formatted_val
                
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
                p.font.name = CorporateTheme.FONT_BODY
                p.font.size = Pt(CorporateTheme.SIZE_BODY - 1)
                p.font.color.rgb = CorporateTheme.PRIMARY
                
                # Apply row backgrounds
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color

        # 4. Set Column Widths evenly across parent width
        col_width = width / num_cols
        for c in range(num_cols):
            table.columns[c].width = col_width
