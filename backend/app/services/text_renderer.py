from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Tuple, Optional
from app.services.theme_manager import CorporateTheme
from app.services.layout_manager import LayoutManager

class TextRenderer:
    """Manages adding title blocks, bullet list points, footers, and font scale calculations."""

    @staticmethod
    def add_slide_header(slide, title_text: str, subtitle_text: Optional[str] = None):
        """Adds a standardized slide header block with correct styling."""
        left, top, width, height = LayoutManager.get_slide_header_box()
        
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Add Title paragraph
        p = tf.paragraphs[0]
        p.text = str(title_text).strip()
        p.font.name = CorporateTheme.FONT_TITLE
        p.font.size = Pt(CorporateTheme.SIZE_TITLE_SLIDE)
        p.font.bold = True
        p.font.color.rgb = CorporateTheme.PRIMARY
        
        # Optional Subtitle
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = str(subtitle_text).strip()
            p2.font.name = CorporateTheme.FONT_BODY
            p2.font.size = Pt(CorporateTheme.SIZE_BODY - 2)
            p2.font.color.rgb = CorporateTheme.SECONDARY
            p2.space_before = Pt(4)

    @staticmethod
    def add_slide_footer(slide, slide_num: int, total_slides: int):
        """Adds standardized footer branding and pagination elements."""
        left, top, width, height = LayoutManager.get_slide_footer_box()
        
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.text = f"DeckMate Automation  |  Slide {slide_num} of {total_slides}"
        p.font.name = CorporateTheme.FONT_BODY
        p.font.size = Pt(CorporateTheme.SIZE_CAPTION)
        p.font.color.rgb = CorporateTheme.SECONDARY

    @staticmethod
    def add_bullet_list(
        slide,
        items: List[str],
        bounding_box: Tuple[Inches, Inches, Inches, Inches],
        list_title: Optional[str] = None
    ):
        """Renders insight bullet items, automatically scaling font size to avoid layout overflows."""
        left, top, width, height = bounding_box
        
        # Draw background container box for side panels (like Key Analytical Insights)
        is_side_panel = list_title in ["Key Analytical Insights", "Operational Recommendations", "ANALYSIS & KEY INSIGHTS"]
        
        if is_side_panel:
            # Draw rounded rectangle background shape
            panel_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
            )
            panel_shape.fill.solid()
            panel_shape.fill.fore_color.rgb = CorporateTheme.BACKGROUND_LIGHT
            panel_shape.line.color.rgb = CorporateTheme.LIGHT_GRAY
            panel_shape.line.width = Pt(1)
            
            # Apply padding margin inset to the text box
            padding = Inches(0.2)
            left_box = left + padding
            top_box = top + padding
            width_box = width - (2 * padding)
            height_box = height - (2 * padding)
        else:
            left_box, top_box, width_box, height_box = left, top, width, height
            
        tx_box = slide.shapes.add_textbox(left_box, top_box, width_box, height_box)
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Dynamic Font Size Auto-Fit
        total_chars = sum(len(x) for x in items) + (len(list_title) if list_title else 0)
        if total_chars > 500:
            font_size = CorporateTheme.SIZE_BODY - 3
        elif total_chars > 300:
            font_size = CorporateTheme.SIZE_BODY - 1.5
        else:
            font_size = CorporateTheme.SIZE_BODY
            
        p_idx = 0
        if list_title:
            p = tf.paragraphs[0]
            p.text = str(list_title).strip()
            p.font.name = CorporateTheme.FONT_BODY
            p.font.size = Pt(font_size + 1.5)
            p.font.bold = True
            p.font.color.rgb = CorporateTheme.PRIMARY
            p.space_after = Pt(8)
            p_idx += 1
            
        for item in items:
            p = tf.add_paragraph() if p_idx > 0 else tf.paragraphs[0]
            item_str = str(item).strip()
            
            p.font.name = CorporateTheme.FONT_BODY
            p.font.size = Pt(font_size)
            p.level = 0
            p.space_after = Pt(6)
            
            # Format bullet indicator
            text_to_format = item_str
            is_bold_prefix = False
            prefix = ""
            suffix = text_to_format
            
            if "**" in text_to_format:
                parts = text_to_format.split("**", 2)
                if len(parts) >= 3:
                    prefix = parts[1]
                    suffix = parts[2]
                    is_bold_prefix = True
            elif ":" in text_to_format:
                parts = text_to_format.split(":", 1)
                prefix = parts[0]
                suffix = ":" + parts[1]
                is_bold_prefix = True
                
            if is_bold_prefix:
                # Add bold run for prefix label
                r1 = p.add_run()
                r1.text = "•  " + prefix.strip()
                r1.font.bold = True
                r1.font.name = CorporateTheme.FONT_BODY
                r1.font.size = Pt(font_size)
                r1.font.color.rgb = CorporateTheme.PRIMARY
                
                # Add normal run for remaining content
                r2 = p.add_run()
                r2.text = suffix
                r2.font.bold = False
                r2.font.name = CorporateTheme.FONT_BODY
                r2.font.size = Pt(font_size)
                r2.font.color.rgb = CorporateTheme.SECONDARY
            else:
                p.text = "•  " + text_to_format
                p.font.color.rgb = CorporateTheme.SECONDARY
                
            p_idx += 1
