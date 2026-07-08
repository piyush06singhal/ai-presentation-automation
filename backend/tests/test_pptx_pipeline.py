import pytest
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from app.schemas.models import (
    StoryboardRequest, 
    SlidePlan, 
    BusinessSummary, 
    WorkbookMetadata, 
    WorksheetMetadata, 
    ColumnSchema, 
    WorkbookHealth,
    DataTypeEnum, 
    DetectedKPI, 
    BusinessTrend, 
    BusinessFact, 
    TrendDirectionEnum
)
from app.services.ppt_compiler import PPTCompiler
from app.services.theme_manager import CorporateTheme
from app.services.export_service import ExportService

@pytest.fixture
def mock_bi_payload():
    """Generates a dummy BusinessSummary and active dataframe collection mapping."""
    # Dataframe setup
    data = {
        "Date": ["2026-01-01", "2026-02-01", "2026-03-01"],
        "Sales": [150, 200, 180],
        "OperatingCost": [40, 42, 38],
        "Category": ["A", "B", "C"]
    }
    df = pd.DataFrame(data)
    df["Date"] = pd.to_datetime(df["Date"])
    df_collection = {"MonthlyPerformance": df}

    # Metadata and summary schema setups
    sheets = [
        WorksheetMetadata(
            name="MonthlyPerformance",
            columns=[
                ColumnSchema(name="Date", datatype=DataTypeEnum.DATE),
                ColumnSchema(name="Sales", datatype=DataTypeEnum.CURRENCY),
                ColumnSchema(name="OperatingCost", datatype=DataTypeEnum.CURRENCY),
                ColumnSchema(name="Category", datatype=DataTypeEnum.CATEGORICAL)
            ],
            row_count=3,
            col_count=4,
            is_empty=False
        )
    ]
    metadata = WorkbookMetadata(sheets=sheets, file_size_bytes=2048, active_sheets_count=1)
    health = WorkbookHealth(issues=[], overall_score=100.0)
    kpis = [
        DetectedKPI(
            name="Sales",
            value=530.0,
            column="Sales",
            worksheet="MonthlyPerformance",
            confidence=1.0,
            unit="$",
            description="Overall sales total."
        )
    ]
    trends = [
        BusinessTrend(
            metric="Sales",
            direction=TrendDirectionEnum.UPWARD,
            percentage_change=33.3,
            time_col="Date",
            growth_rate=0.2,
            description="Upward sales growth",
            worksheet="MonthlyPerformance"
        )
    ]
    facts = [
        BusinessFact(
            fact_id="FACT_001",
            statement="Sales climbed from $150 to $200 in Month 2.",
            source_worksheet="MonthlyPerformance",
            metrics=["Sales"],
            importance=3
        )
    ]
    
    summary = BusinessSummary(
        metadata=metadata,
        health=health,
        kpis=kpis,
        trends=trends,
        facts=facts,
        statistics={"MonthlyPerformance": {}}
    )
    
    return summary, df_collection

@pytest.fixture
def mock_presentation_plan():
    """Builds a diverse 6-slide narrative presentation plan matching different layouts."""
    slides = [
        SlidePlan(
            slide_id="slide_1",
            template_id="Title",
            title="Q3 Business Performance Review",
            objective="Analyze core operations and revenue milestones.",
            worksheet="MonthlyPerformance",
            insights=[],
            speaker_notes="Welcome all to the quarterly meeting.",
            why_created="Introduce deck structure",
            priority=1,
            confidence=1.0
        ),
        SlidePlan(
            slide_id="slide_2",
            template_id="Agenda",
            title="Agenda & Key Data Sections",
            objective="Workbook structures list",
            worksheet="MonthlyPerformance",
            insights=[],
            speaker_notes="Overview of worksheets",
            why_created="Establish schedule",
            priority=1,
            confidence=1.0
        ),
        SlidePlan(
            slide_id="slide_3",
            template_id="KPI Dashboard",
            title="Core Operational KPIs",
            objective="High-level metrics display",
            worksheet="MonthlyPerformance",
            insights=["Sales value shows steady volume support.", "Cost structures remain contained."],
            speaker_notes="KPI card breakdown",
            why_created="Summary metrics highlight",
            priority=1,
            confidence=1.0
        ),
        SlidePlan(
            slide_id="slide_4",
            template_id="Trend Analysis",
            title="Sales Inflow Analysis",
            objective="Date trends overview",
            worksheet="MonthlyPerformance",
            chart_type="Line",
            x_axis="Date",
            y_axis=["Sales", "OperatingCost"],
            insights=["Consistent growth month-over-month.", "Seasonal trend support noted."],
            speaker_notes="Speaker notes for trend lines.",
            why_created="Show trends",
            priority=2,
            confidence=0.9
        ),
        SlidePlan(
            slide_id="slide_5",
            template_id="Recommendations",
            title="Strategic Interventions",
            objective="Growth opportunities highlights",
            worksheet="MonthlyPerformance",
            insights=["General cost controls recommended."],
            recommendations=["Scale product marketing assets.", "Perform detailed segment review."],
            speaker_notes="Operational checklist",
            why_created="Next steps proposal",
            priority=1,
            confidence=0.95
        ),
        SlidePlan(
            slide_id="slide_6",
            template_id="Appendix",
            title="Appendix & Data Checks",
            objective="Visual details audit logs",
            worksheet="MonthlyPerformance",
            insights=["Workbook audit completed without warning blocks."],
            speaker_notes="Final details notes",
            why_created="Logs log",
            priority=3,
            confidence=1.0
        )
    ]
    return StoryboardRequest(
        audience="CEO",
        objective="Analyze revenue performance",
        slides=slides
    )

def test_widescreen_ppt_compilation(mock_presentation_plan, mock_bi_payload):
    summary, df_collection = mock_bi_payload
    ppt_bytes = PPTCompiler.compile_presentation(
        plan=mock_presentation_plan,
        summary=summary,
        df_collection=df_collection
    )
    
    # Verify stream size
    assert len(ppt_bytes) > 0
    
    # Reload stream using pptx and verify slide parameters
    stream = io.BytesIO(ppt_bytes)
    prs = Presentation(stream)
    
    # Verify widescreen sizes
    assert prs.slide_width == CorporateTheme.SLIDE_WIDTH
    assert prs.slide_height == CorporateTheme.SLIDE_HEIGHT
    
    # Verify slide counts matching
    assert len(prs.slides) == 6
    
    # Verify Title slide exists
    title_slide = prs.slides[0]
    # Blank layouts shouldn't have default placeholders but custom drawn shapes
    assert len(title_slide.shapes) >= 2

def test_chart_to_table_fallback(mock_bi_payload):
    summary, df_collection = mock_bi_payload
    
    # Request a slide containing an invalid chart column mapping
    slide_invalid = SlidePlan(
        slide_id="slide_invalid",
        template_id="TrendAnalysis",
        title="Invalid Columns Trend",
        objective="Validation Check",
        worksheet="MonthlyPerformance",
        chart_type="Line",
        x_axis="Non_Existent_X_Column", # Will fail chart rendering
        y_axis=["Sales"],
        insights=["Fallback tests"],
        speaker_notes="Notes",
        why_created="Check fallback behavior",
        priority=1,
        confidence=1.0
    )
    
    plan = StoryboardRequest(
        audience="CEO",
        objective="Test fallbacks",
        slides=[slide_invalid]
    )
    
    ppt_bytes = PPTCompiler.compile_presentation(
        plan=plan,
        summary=summary,
        df_collection=df_collection
    )
    
    assert len(ppt_bytes) > 0
    prs = Presentation(io.BytesIO(ppt_bytes))
    
    # The slide compiles successfully without raising exceptions (uses visual warning or table layout)
    assert len(prs.slides) == 1

def test_export_service_helpers(mock_presentation_plan, mock_bi_payload):
    summary, df_collection = mock_bi_payload
    ppt_bytes = PPTCompiler.compile_presentation(
        plan=mock_presentation_plan,
        summary=summary,
        df_collection=df_collection
    )
    
    # Verify stream wrapping
    stream = ExportService.export_to_stream(ppt_bytes)
    assert isinstance(stream, io.BytesIO)
    assert stream.getvalue() == ppt_bytes
